"""PowerPoint (.pptx) document parser using python-pptx."""

import logging

from pptx import Presentation

logger = logging.getLogger(__name__)


def parse_pptx(file_path: str, source: str = "") -> dict:
    """Extract text from a PowerPoint file.

    Each slide is rendered with a slide number heading, followed by
    the text from all shapes.
    Returns: {"text": str, "metadata": dict}
    """
    prs = Presentation(file_path)
    slides_text = []

    for i, slide in enumerate(prs.slides, 1):
        parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if text:
                    parts.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        parts.append(" | ".join(cells))
        if parts:
            slides_text.append(f"## Slide {i}\n\n" + "\n\n".join(parts))

    full_text = "\n\n".join(slides_text)

    return {
        "text": full_text,
        "metadata": {
            "source": source,
            "format": "pptx",
            "slide_count": len(prs.slides),
        },
    }
